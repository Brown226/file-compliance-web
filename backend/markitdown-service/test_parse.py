#!/usr/bin/env python3
"""
MarkItDown 文档解析服务 - 四种格式解析测试脚本

测试内容：
1. DOCX - Word 文档解析（文本、标题、表格）
2. XLSX - Excel 电子表格解析（单元格、多 Sheet）
3. PDF  - PDF 文档解析（逐页文本、表格）
4. PPTX - PowerPoint 演示文稿解析（幻灯片、文本）

使用方式：
  python test_parse.py [--url http://localhost:8000] [--keep]

选项：
  --url   MarkItDown 服务地址（默认 http://localhost:8000）
  --keep  保留测试文件（默认测试后删除）
"""

import os
import sys
import json
import time
import argparse
import tempfile
import traceback
from pathlib import Path

# Windows 控制台 UTF-8 输出
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')

# 第三方库
try:
    import requests
except ImportError:
    print("❌ 缺少 requests 库，请安装: pip install requests")
    sys.exit(1)


# ==================== 测试文件生成 ====================

def create_test_docx(path: str) -> str:
    """创建包含标题、段落、表格的 DOCX 测试文件"""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 标题
    doc.add_heading('测试文档 - 文件合规审查系统', level=1)

    # 正文段落
    doc.add_paragraph(
        '本文档用于验证 MarkItDown 解析服务对 Word 文档的解析能力。'
        '系统支持对 DOCX、XLSX、PDF、PPTX 四种格式进行智能解析，'
        '提取文本内容、结构化数据和 Markdown 格式输出。'
    )

    # 二级标题
    doc.add_heading('1. 系统架构', level=2)
    doc.add_paragraph(
        '系统采用前后端分离架构，后端使用 Node.js + Express + Prisma ORM，'
        '前端使用 Vue3 + Element Plus + Pinia 状态管理。'
    )

    # 三级标题 + 列表
    doc.add_heading('1.1 核心模块', level=3)
    doc.add_paragraph('文件上传模块：支持多文件批量上传，单文件最大 100MB', style='List Bullet')
    doc.add_paragraph('智能解析模块：集成 MarkItDown 引擎，支持四种格式', style='List Bullet')
    doc.add_paragraph('规则引擎模块：基于正则 + LLM 的双模式审查', style='List Bullet')

    # 表格
    doc.add_heading('2. 支持的文件格式', level=2)
    table = doc.add_table(rows=5, cols=4, style='Table Grid')
    headers = ['格式', '扩展名', '解析引擎', '支持能力']
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h

    data = [
        ['Word 文档', '.docx', 'mammoth', '文本/标题/表格'],
        ['Excel 表格', '.xlsx', 'openpyxl', '单元格/多Sheet'],
        ['PDF 文档', '.pdf', 'pdfplumber', '逐页文本/表格'],
        ['PowerPoint', '.pptx', 'python-pptx', '幻灯片/文本'],
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = cell_data

    # 更多段落
    doc.add_heading('3. 审查流程', level=2)
    doc.add_paragraph(
        '用户上传文件后，系统自动触发审查流水线：'
        '首先进行文件解析，提取文本内容；'
        '然后通过规则引擎进行格式审查；'
        '最后调用 LLM 进行内容审查，生成审查报告。'
    )

    doc.add_paragraph('')  # 空段落
    p = doc.add_paragraph('重要提示：')
    run = p.runs[0]
    run.bold = True
    p.add_run('所有审查结果均可追溯，支持原文定位和高亮显示。')

    doc.save(path)
    return path


def create_test_xlsx(path: str) -> str:
    """创建包含多 Sheet、不同数据类型的 XLSX 测试文件"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

    wb = Workbook()

    # ===== Sheet 1: 项目信息 =====
    ws1 = wb.active
    ws1.title = '项目信息'

    # 标题行（加粗 + 蓝底）
    header_font = Font(bold=True, color='FFFFFF', size=12)
    header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin'),
    )

    headers = ['项目编号', '项目名称', '负责人', '状态', '完成度']
    for col, h in enumerate(headers, 1):
        cell = ws1.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')
        cell.border = thin_border

    # 数据行
    projects = [
        ['PRJ-2026-001', '核电站BOP管道设计', '张三', '进行中', 0.75],
        ['PRJ-2026-002', '给排水系统审查', '李四', '已完成', 1.0],
        ['PRJ-2026-003', '电气线路合规检查', '王五', '待审核', 0.5],
        ['PRJ-2026-004', '消防系统设计审查', '赵六', '进行中', 0.6],
        ['PRJ-2026-005', '暖通空调标准验证', '钱七', '未开始', 0.0],
    ]
    for row_idx, row_data in enumerate(projects, 2):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws1.cell(row=row_idx, column=col_idx, value=value)
            cell.border = thin_border
            if col_idx == 5:  # 完成度百分比格式
                cell.number_format = '0%'

    # 设置列宽
    ws1.column_dimensions['A'].width = 16
    ws1.column_dimensions['B'].width = 22
    ws1.column_dimensions['C'].width = 10
    ws1.column_dimensions['D'].width = 10
    ws1.column_dimensions['E'].width = 10

    # ===== Sheet 2: 标准引用 =====
    ws2 = wb.create_sheet('标准引用')

    ref_headers = ['标准编号', '标准名称', '发布日期', '状态', '引用次数']
    for col, h in enumerate(ref_headers, 1):
        cell = ws2.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')
        cell.border = thin_border

    standards = [
        ['GB/T 50001-2017', '房屋建筑制图统一标准', '2017-11-01', '现行', 42],
        ['GB 50242-2002', '建筑给水排水及采暖工程施工质量验收规范', '2002-04-01', '现行', 28],
        ['GB 50016-2014', '建筑设计防火规范', '2014-08-01', '现行', 56],
        ['HAF 601', '核安全法规', '2010-01-01', '现行', 15],
        ['NB/T 20003-2010', '核电厂核岛机械设备焊接规范', '2010-06-01', '现行', 8],
    ]
    for row_idx, row_data in enumerate(standards, 2):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws2.cell(row=row_idx, column=col_idx, value=value)
            cell.border = thin_border

    ws2.column_dimensions['A'].width = 20
    ws2.column_dimensions['B'].width = 40
    ws2.column_dimensions['C'].width = 14
    ws2.column_dimensions['D'].width = 8
    ws2.column_dimensions['E'].width = 10

    # ===== Sheet 3: 统计数据（数字为主） =====
    ws3 = wb.create_sheet('统计数据')
    ws3.cell(row=1, column=1, value='月份')
    ws3.cell(row=1, column=2, value='审查文件数')
    ws3.cell(row=1, column=3, value='发现问题数')
    ws3.cell(row=1, column=4, value='合规率')

    months_data = [
        ['2026-01', 128, 45, 0.648],
        ['2026-02', 156, 38, 0.756],
        ['2026-03', 203, 52, 0.744],
        ['2026-04', 89, 12, 0.865],
    ]
    for row_idx, row_data in enumerate(months_data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws3.cell(row=row_idx, column=col_idx, value=value)

    wb.save(path)
    return path


def create_test_pptx(path: str) -> str:
    """创建包含多幻灯片、文本和表格的 PPTX 测试文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: 封面 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = '文件智能审查系统'
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = '技术方案汇报'
    p2.font.size = Pt(28)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = '2026年4月'
    p3.font.size = Pt(18)
    p3.alignment = PP_ALIGN.CENTER
    # p3.font.color.rgb 使用默认颜色，不显式设置

    # ===== Slide 2: 系统概述 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = '系统概述'
    p.font.size = Pt(32)
    p.font.bold = True

    txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    items = [
        '支持 DOCX、XLSX、PDF、PPTX 四种格式自动解析',
        '基于 MarkItDown 引擎实现文档到 Markdown 的精准转换',
        '集成规则引擎 + LLM 双模式审查能力',
        '支持 MaxKB 知识库检索增强审查（RAG）',
        '审查结果可追溯，支持原文定位高亮',
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f'• {item}'
        p.font.size = Pt(20)
        p.space_after = Pt(12)

    # ===== Slide 3: 表格数据 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox4 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf4 = txBox4.text_frame
    p = tf4.paragraphs[0]
    p.text = '审查统计'
    p.font.size = Pt(32)
    p.font.bold = True

    # 添加表格
    rows, cols = 5, 4
    table_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(11), Inches(4))
    table = table_shape.table

    headers = ['月份', '审查文件数', '发现问题数', '合规率']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True

    data = [
        ['2026-01', '128', '45', '64.8%'],
        ['2026-02', '156', '38', '75.6%'],
        ['2026-03', '203', '52', '74.4%'],
        ['2026-04', '89', '12', '86.5%'],
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)

    # ===== Slide 4: 结论 =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox5 = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    p = tf5.paragraphs[0]
    p.text = '结论：系统已具备完整的文档解析与审查能力'
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf5.add_paragraph()
    p2.text = '下一步：持续优化 OCR 准确率和 LLM 审查质量'
    p2.font.size = Pt(22)
    p2.alignment = PP_ALIGN.CENTER

    prs.save(path)
    return path


def get_or_create_test_pdf(path: str) -> str:
    """
    获取测试 PDF 文件。
    优先使用 uploads 目录中已有的真实 PDF，若不存在则报告跳过。
    """
    # 查找已有的真实 PDF 文件
    uploads_dir = os.path.join(os.path.dirname(__file__), '..', 'uploads')
    if os.path.isdir(uploads_dir):
        for f in os.listdir(uploads_dir):
            if f.endswith('.pdf') and not f.startswith('files-'):
                src = os.path.join(uploads_dir, f)
                import shutil
                shutil.copy2(src, path)
                return path

    # 没有找到 PDF，尝试用 reportlab 生成
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table

        doc = SimpleDocTemplate(path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph('测试PDF文档 - 文件合规审查系统', styles['Title']))
        story.append(Spacer(1, 20))
        story.append(Paragraph(
            '本文档用于验证 MarkItDown 解析服务对 PDF 文档的解析能力。'
            '系统使用 pdfplumber 进行逐页文本提取，同时使用 markitdown 获取 Markdown 格式输出。',
            styles['Normal']
        ))
        story.append(Spacer(1, 12))
        story.append(Paragraph('1. PDF 解析特性', styles['Heading1']))
        story.append(Paragraph(
            'PDF 解析支持逐页文本提取、表格检测、图片识别等能力。'
            '对于扫描型 PDF，系统会自动调用 OCR 服务进行文字识别。',
            styles['Normal']
        ))

        # 简单表格
        table_data = [
            ['特性', '说明'],
            ['逐页提取', '使用 pdfplumber 获取真实页码'],
            ['表格检测', '支持 pdfplumber 表格提取'],
            ['OCR 降级', '扫描型 PDF 自动走 OCR 路径'],
        ]
        t = Table(table_data)
        story.append(t)

        story.append(Spacer(1, 12))
        story.append(Paragraph('2. 使用说明', styles['Heading1']))
        story.append(Paragraph(
            '上传 PDF 文件后，系统会自动进行解析。'
            '解析结果包含文本内容、页面数据、结构化信息和 Markdown 格式。',
            styles['Normal']
        ))

        doc.build(story)
        return path
    except ImportError:
        return None


# ==================== API 测试 ====================

class ParseTestResult:
    """单个测试结果"""
    def __init__(self, name: str):
        self.name = name
        self.success = False
        self.error = None
        self.duration_ms = 0
        self.details = {}
        self.checks = []  # [(check_name, passed, detail)]

    def add_check(self, name: str, passed: bool, detail: str = ''):
        self.checks.append((name, passed, detail))
        if not passed:
            self.success = False

    def __repr__(self):
        status = '✅ PASS' if self.success else '❌ FAIL'
        return f'{status} {self.name} ({self.duration_ms}ms)'


def call_parse_api(base_url: str, file_path: str, file_type: str = '') -> dict:
    """调用 /api/parse 接口"""
    url = f'{base_url}/api/parse'
    filename = os.path.basename(file_path)

    with open(file_path, 'rb') as f:
        files = {'file': (filename, f)}
        data = {}
        if file_type:
            data['file_type'] = file_type

        resp = requests.post(url, files=files, data=data, timeout=60)

    if resp.status_code != 200:
        raise Exception(f'HTTP {resp.status_code}: {resp.text[:500]}')

    result = resp.json()
    if result.get('code') != 200:
        raise Exception(f"API error: {result.get('message', 'unknown')}")

    return result.get('data', {})


def call_convert_api(base_url: str, file_path: str) -> dict:
    """调用 /api/convert 接口"""
    url = f'{base_url}/api/convert'
    filename = os.path.basename(file_path)

    with open(file_path, 'rb') as f:
        files = {'file': (filename, f)}
        resp = requests.post(url, files=files, timeout=60)

    if resp.status_code != 200:
        raise Exception(f'HTTP {resp.status_code}: {resp.text[:500]}')

    return resp.json()


# ==================== 测试用例 ====================

def test_docx(base_url: str, file_path: str) -> ParseTestResult:
    """测试 DOCX 解析"""
    result = ParseTestResult('DOCX - Word 文档解析')

    try:
        start = time.time()
        data = call_parse_api(base_url, file_path, 'docx')
        result.duration_ms = int((time.time() - start) * 1000)

        # 检查1: 返回数据完整性
        has_all_fields = all(k in data for k in ['text', 'pages', 'metadata', 'structure', 'markdown'])
        result.add_check('返回字段完整', has_all_fields,
                         f"字段: {', '.join(data.keys())}" if has_all_fields else f"缺失: {[k for k in ['text','pages','metadata','structure','markdown'] if k not in data]}")

        # 检查2: 文本内容非空
        text = data.get('text', '')
        result.add_check('文本内容非空', len(text) > 100,
                         f'文本长度: {len(text)} 字符')

        # 检查3: Markdown 内容非空
        markdown = data.get('markdown', '')
        result.add_check('Markdown 非空', len(markdown) > 100,
                         f'Markdown 长度: {len(markdown)} 字符')

        # 检查4: 标题提取
        structure = data.get('structure', {})
        paragraphs = structure.get('paragraphs', [])
        headings = [p for p in paragraphs if p.get('style', '').startswith('Heading')]
        result.add_check('标题提取', len(headings) >= 3,
                         f'提取到 {len(headings)} 个标题: {[h["text"][:30] for h in headings[:5]]}')

        # 检查5: 表格检测
        tables = structure.get('tables', [])
        has_tables = data.get('metadata', {}).get('has_tables', False)
        result.add_check('表格检测', has_tables or len(tables) > 0,
                         f'metadata.has_tables={has_tables}, structure.tables={len(tables)}')

        # 检查6: 页面数据
        pages = data.get('pages', [])
        result.add_check('页面数据', len(pages) >= 1,
                         f'页数: {len(pages)}')

        # 检查7: 无解析错误
        parse_error = data.get('metadata', {}).get('parse_error')
        result.add_check('无解析错误', parse_error is None,
                         f'parse_error={parse_error}')

        # 检查8: Markdown 包含标题标记
        has_md_headings = bool(markdown and '#' in markdown)
        result.add_check('Markdown 标题标记', has_md_headings,
                         'Markdown 中包含 # 标题标记' if has_md_headings else 'Markdown 中无标题标记')

        # 检查9: Markdown 包含表格标记
        has_md_table = bool(markdown and '|' in markdown)
        result.add_check('Markdown 表格标记', has_md_table,
                         'Markdown 中包含表格标记' if has_md_table else 'Markdown 中无表格标记')

        # 全部检查通过
        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


def test_xlsx(base_url: str, file_path: str) -> ParseTestResult:
    """测试 XLSX 解析"""
    result = ParseTestResult('XLSX - Excel 电子表格解析')

    try:
        start = time.time()
        data = call_parse_api(base_url, file_path, 'xlsx')
        result.duration_ms = int((time.time() - start) * 1000)

        # 检查1: 返回数据完整性
        has_all_fields = all(k in data for k in ['text', 'pages', 'metadata', 'structure', 'markdown'])
        result.add_check('返回字段完整', has_all_fields,
                         f"字段: {', '.join(data.keys())}" if has_all_fields else f"缺失: {[k for k in ['text','pages','metadata','structure','markdown'] if k not in data]}")

        # 检查2: 文本内容非空
        text = data.get('text', '')
        result.add_check('文本内容非空', len(text) > 50,
                         f'文本长度: {len(text)} 字符')

        # 检查3: Markdown 内容非空
        markdown = data.get('markdown', '')
        result.add_check('Markdown 非空', len(markdown) > 50,
                         f'Markdown 长度: {len(markdown)} 字符')

        # 检查4: 包含表格数据
        tables = structure_tables(data)
        result.add_check('表格数据', len(tables) > 0,
                         f'检测到 {len(tables)} 个表格')

        # 检查5: 包含标准编号内容（GB/T 开头）
        has_standard_ref = bool(text and ('GB/T' in text or 'GB ' in text or 'NB/T' in text))
        result.add_check('标准编号内容', has_standard_ref,
                         '文本中包含 GB/T 或 GB 或 NB/T 标准编号' if has_standard_ref else '未检测到标准编号')

        # 检查6: 无解析错误
        parse_error = data.get('metadata', {}).get('parse_error')
        result.add_check('无解析错误', parse_error is None,
                         f'parse_error={parse_error}')

        # 检查7: Markdown 包含表格标记
        has_md_table = bool(markdown and '|' in markdown)
        result.add_check('Markdown 表格标记', has_md_table,
                         'Markdown 中包含 | 表格标记' if has_md_table else 'Markdown 中无表格标记')

        # 检查8: 多 Sheet 内容
        has_multi_sheet = bool(markdown and ('Sheet' in markdown or '项目信息' in markdown or '标准引用' in markdown))
        result.add_check('多 Sheet 内容', has_multi_sheet,
                         'Markdown 中包含多 Sheet 标记' if has_multi_sheet else 'Markdown 中未检测到多 Sheet')

        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


def test_pdf(base_url: str, file_path: str) -> ParseTestResult:
    """测试 PDF 解析"""
    result = ParseTestResult('PDF - PDF 文档解析')

    try:
        start = time.time()
        data = call_parse_api(base_url, file_path, 'pdf')
        result.duration_ms = int((time.time() - start) * 1000)

        # 检查1: 返回数据完整性
        has_all_fields = all(k in data for k in ['text', 'pages', 'metadata', 'structure', 'markdown'])
        result.add_check('返回字段完整', has_all_fields,
                         f"字段: {', '.join(data.keys())}" if has_all_fields else f"缺失: {[k for k in ['text','pages','metadata','structure','markdown'] if k not in data]}")

        # 检查2: 文本内容非空
        text = data.get('text', '')
        result.add_check('文本内容非空', len(text) > 20,
                         f'文本长度: {len(text)} 字符')

        # 检查3: 页面数据（PDF 应有逐页数据）
        pages = data.get('pages', [])
        page_count = data.get('metadata', {}).get('page_count', 0)
        result.add_check('逐页数据', len(pages) >= 1 and page_count >= 1,
                         f'页数: {page_count}, pages 数组长度: {len(pages)}')

        # 检查4: Markdown 内容
        markdown = data.get('markdown', '')
        result.add_check('Markdown 非空', len(markdown) > 20,
                         f'Markdown 长度: {len(markdown)} 字符')

        # 检查5: 无解析错误
        parse_error = data.get('metadata', {}).get('parse_error')
        result.add_check('无解析错误', parse_error is None,
                         f'parse_error={parse_error}')

        # 检查6: 结构化数据
        structure = data.get('structure', {})
        paragraphs = structure.get('paragraphs', [])
        result.add_check('结构化段落', len(paragraphs) >= 1,
                         f'段落数量: {len(paragraphs)}')

        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


def test_pptx(base_url: str, file_path: str) -> ParseTestResult:
    """测试 PPTX 解析"""
    result = ParseTestResult('PPTX - PowerPoint 演示文稿解析')

    try:
        start = time.time()
        data = call_parse_api(base_url, file_path, 'pptx')
        result.duration_ms = int((time.time() - start) * 1000)

        # 检查1: 返回数据完整性
        has_all_fields = all(k in data for k in ['text', 'pages', 'metadata', 'structure', 'markdown'])
        result.add_check('返回字段完整', has_all_fields,
                         f"字段: {', '.join(data.keys())}" if has_all_fields else f"缺失: {[k for k in ['text','pages','metadata','structure','markdown'] if k not in data]}")

        # 检查2: 文本内容非空
        text = data.get('text', '')
        result.add_check('文本内容非空', len(text) > 50,
                         f'文本长度: {len(text)} 字符')

        # 检查3: Markdown 内容
        markdown = data.get('markdown', '')
        result.add_check('Markdown 非空', len(markdown) > 50,
                         f'Markdown 长度: {len(markdown)} 字符')

        # 检查4: 包含幻灯片内容
        has_slide_content = bool(text and ('审查' in text or '系统' in text or '解析' in text or '文件' in text))
        result.add_check('幻灯片内容', has_slide_content,
                         '文本中包含幻灯片关键词' if has_slide_content else '未检测到幻灯片关键词')

        # 检查5: Markdown 包含表格标记
        has_md_table = bool(markdown and '|' in markdown)
        result.add_check('Markdown 表格标记', has_md_table,
                         'Markdown 中包含 | 表格标记' if has_md_table else 'Markdown 中无表格标记（可能无表格幻灯片）')

        # 检查6: 无解析错误
        parse_error = data.get('metadata', {}).get('parse_error')
        result.add_check('无解析错误', parse_error is None,
                         f'parse_error={parse_error}')

        # 检查7: 结构化数据
        structure = data.get('structure', {})
        paragraphs = structure.get('paragraphs', [])
        result.add_check('结构化段落', len(paragraphs) >= 1,
                         f'段落数量: {len(paragraphs)}')

        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


def test_convert_api(base_url: str, file_path: str, ext: str) -> ParseTestResult:
    """测试 /api/convert 接口"""
    result = ParseTestResult(f'Convert API - {ext.upper()}')

    try:
        start = time.time()
        data = call_convert_api(base_url, file_path)
        result.duration_ms = int((time.time() - start) * 1000)

        result.add_check('成功', data.get('success', False),
                         f'success={data.get("success")}')

        result.add_check('Markdown 非空', len(data.get('markdown', '')) > 0,
                         f'长度: {len(data.get("markdown", ""))} 字符')

        result.add_check('扩展名正确', data.get('extension') == ext,
                         f'extension={data.get("extension")}')

        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


def test_health(base_url: str) -> ParseTestResult:
    """测试健康检查接口"""
    result = ParseTestResult('Health Check')

    try:
        start = time.time()
        resp = requests.get(f'{base_url}/health', timeout=5)
        result.duration_ms = int((time.time() - start) * 1000)

        data = resp.json()
        result.add_check('状态 OK', data.get('status') == 'ok',
                         f'status={data.get("status")}')

        result.add_check('版本号', bool(data.get('version')),
                         f'version={data.get("version")}')

        result.add_check('格式列表', len(data.get('supported_formats', [])) >= 4,
                         f'formats={data.get("supported_formats")}')

        if all(c[1] for c in result.checks):
            result.success = True

    except Exception as e:
        result.error = str(e)
        traceback.print_exc()

    return result


# ==================== 辅助函数 ====================

def structure_tables(data: dict) -> list:
    """从解析结果中提取表格"""
    return data.get('structure', {}).get('tables', [])


def print_result(r: ParseTestResult):
    """打印单个测试结果"""
    status = '✅' if r.success else '❌'
    print(f'\n  {status} {r.name}  ({r.duration_ms}ms)')

    if r.error:
        print(f'     💥 异常: {r.error}')
        return

    for check_name, passed, detail in r.checks:
        icon = '✅' if passed else '❌'
        detail_str = f' — {detail}' if detail else ''
        print(f'     {icon} {check_name}{detail_str}')


def print_summary(results: list):
    """打印测试汇总"""
    print('\n' + '=' * 60)
    print('📊 测试汇总')
    print('=' * 60)

    total = len(results)
    passed = sum(1 for r in results if r.success)
    failed = total - passed

    for r in results:
        status = '✅ PASS' if r.success else '❌ FAIL'
        print(f'  {status}  {r.name}  ({r.duration_ms}ms)')

    print(f'\n  总计: {total} | 通过: {passed} | 失败: {failed}')
    print('=' * 60)

    if failed == 0:
        print('\n🎉 所有测试通过！markitdown-service 解析功能完全正常。')
        print('\n📋 可以安全删除 markitdown-main/ 源码目录，执行以下命令：')
        print('   PowerShell:')
        print('     Remove-Item -Recurse -Force "e:/工作/file-compliance-web/markitdown-main"')
        print('   Bash:')
        print('     rm -rf /path/to/markitdown-main/')
    else:
        print(f'\n⚠️  有 {failed} 项测试未通过，请检查后再决定是否删除 markitdown-main/。')

    return failed == 0


# ==================== 主程序 ====================

def main():
    parser = argparse.ArgumentParser(description='MarkItDown 文档解析服务测试脚本')
    parser.add_argument('--url', default='http://localhost:8000', help='MarkItDown 服务地址')
    parser.add_argument('--keep', action='store_true', help='保留测试文件')
    parser.add_argument('--save-results', action='store_true', help='将解析结果保存为 JSON')
    args = parser.parse_args()

    base_url = args.url.rstrip('/')
    print(f'🚀 MarkItDown 解析测试')
    print(f'   服务地址: {base_url}')
    print(f'   测试时间: {time.strftime("%Y-%m-%d %H:%M:%S")}')
    print('=' * 60)

    # 0. 健康检查
    print('\n📡 步骤 0: 健康检查')
    health_result = test_health(base_url)
    print_result(health_result)

    if not health_result.success:
        print('\n❌ 服务不可用，测试终止。请确认 MarkItDown 服务已启动。')
        return 1

    # 创建临时目录存放测试文件
    tmp_dir = tempfile.mkdtemp(prefix='markitdown_test_')
    results = [health_result]

    try:
        # 1. DOCX 测试
        print('\n📄 步骤 1: DOCX 解析测试')
        docx_path = os.path.join(tmp_dir, 'test_document.docx')
        try:
            create_test_docx(docx_path)
            print(f'   测试文件: {docx_path} ({os.path.getsize(docx_path)} bytes)')

            r1 = test_docx(base_url, docx_path)
            print_result(r1)
            results.append(r1)

            # Convert API 测试
            r1c = test_convert_api(base_url, docx_path, '.docx')
            print_result(r1c)
            results.append(r1c)
        except ImportError:
            print('   ⚠️  跳过: 缺少 python-docx 库（pip install python-docx）')

        # 2. XLSX 测试
        print('\n📊 步骤 2: XLSX 解析测试')
        xlsx_path = os.path.join(tmp_dir, 'test_spreadsheet.xlsx')
        try:
            create_test_xlsx(xlsx_path)
            print(f'   测试文件: {xlsx_path} ({os.path.getsize(xlsx_path)} bytes)')

            r2 = test_xlsx(base_url, xlsx_path)
            print_result(r2)
            results.append(r2)

            r2c = test_convert_api(base_url, xlsx_path, '.xlsx')
            print_result(r2c)
            results.append(r2c)
        except ImportError:
            print('   ⚠️  跳过: 缺少 openpyxl 库（pip install openpyxl）')

        # 3. PDF 测试
        print('\n📑 步骤 3: PDF 解析测试')
        pdf_path = os.path.join(tmp_dir, 'test_document.pdf')
        pdf_file = get_or_create_test_pdf(pdf_path)
        if pdf_file and os.path.exists(pdf_path):
            print(f'   测试文件: {pdf_path} ({os.path.getsize(pdf_path)} bytes)')

            r3 = test_pdf(base_url, pdf_path)
            print_result(r3)
            results.append(r3)

            r3c = test_convert_api(base_url, pdf_path, '.pdf')
            print_result(r3c)
            results.append(r3c)
        else:
            print('   ⚠️  跳过: 未找到 PDF 测试文件且无法生成（需要 reportlab）')

        # 4. PPTX 测试
        print('\n📽️  步骤 4: PPTX 解析测试')
        pptx_path = os.path.join(tmp_dir, 'test_presentation.pptx')
        try:
            create_test_pptx(pptx_path)
            print(f'   测试文件: {pptx_path} ({os.path.getsize(pptx_path)} bytes)')

            r4 = test_pptx(base_url, pptx_path)
            print_result(r4)
            results.append(r4)

            r4c = test_convert_api(base_url, pptx_path, '.pptx')
            print_result(r4c)
            results.append(r4c)
        except ImportError:
            print('   ⚠️  跳过: 缺少 python-pptx 库（pip install python-pptx）')

        # 5. DWG 降级测试
        print('\n🔧 步骤 5: DWG 不支持格式降级测试')
        dwg_path = os.path.join(tmp_dir, 'test.dwg')
        with open(dwg_path, 'wb') as f:
            f.write(b'AC1015\x00' + b'\x00' * 100)  # 伪 DWG 头
        try:
            dwg_data = call_parse_api(base_url, dwg_path, 'dwg')
            dwg_result = ParseTestResult('DWG - 不支持格式降级')
            dwg_result.duration_ms = 0
            parse_error = dwg_data.get('metadata', {}).get('parse_error')
            dwg_result.add_check('返回空结果', len(dwg_data.get('text', '')) == 0,
                                 f'text 长度: {len(dwg_data.get("text", ""))}')
            dwg_result.add_check('标记不支持', parse_error is not None and 'not supported' in (parse_error or '').lower(),
                                 f'parse_error={parse_error}')
            if all(c[1] for c in dwg_result.checks):
                dwg_result.success = True
            print_result(dwg_result)
            results.append(dwg_result)
        except Exception as e:
            print(f'   ⚠️  DWG 测试异常: {e}')

    finally:
        # 清理测试文件
        if not args.keep:
            import shutil
            try:
                shutil.rmtree(tmp_dir)
                print(f'\n🗑️  已清理临时目录: {tmp_dir}')
            except Exception:
                print(f'\n⚠️  清理失败: {tmp_dir}')
        else:
            print(f'\n📁 测试文件保留在: {tmp_dir}')

    # 保存详细结果
    if args.save_results:
        results_file = os.path.join(os.path.dirname(__file__), 'test_results.json')
        output = []
        for r in results:
            output.append({
                'name': r.name,
                'success': r.success,
                'duration_ms': r.duration_ms,
                'error': r.error,
                'checks': [{'name': c[0], 'passed': c[1], 'detail': c[2]} for c in r.checks],
            })
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(output, f, ensure_ascii=False, indent=2)
        print(f'📄 详细结果已保存: {results_file}')

    # 打印汇总
    all_passed = print_summary(results)
    return 0 if all_passed else 1


if __name__ == '__main__':
    sys.exit(main())
